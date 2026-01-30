from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title, content_points, layout_index=1):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Check for sub-points (simple indentation logic)
            if point.startswith("    "):
                p.level = 1
                p.text = point.strip()
            else:
                p.level = 0
                p.text = point
            
            p.font.size = Pt(18)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Microwave Oven Simulator using PIC16F877A"
    subtitle.text = "Embedded Systems Internship Project\nEmertxe Information Technologies"

    # Slide 2: Overview
    add_slide("Overview", [
        "Project Title: Microwave Oven Simulator",
        "Core Concept:",
        "    Simulation of a fully functional microwave oven control system using PIC16F877A.",
        "    Focuses on replicating user interface workflows, timing, and load control.",
        "Key Features:",
        "    Modes: Micro, Grill, Convection, Start.",
        "    Interface: 4x4 Matrix Keypad, 16x4 CLCD.",
        "    Timing: Interrupt-based accurate countdown timer.",
        "    Safety: Door status simulation and safe mode transitions."
    ])

    # Slide 3: Goals of Internship
    add_slide("Goals of Internship", [
        "Primary Objective:",
        "    Bridge gap between academic theory and industry-standard embedded development.",
        "Learning Outcomes:",
        "    Microcontroller Architecture: In-depth understanding of PIC16F877A.",
        "    Embedded C Programming: Mastery of bitwise operations and register manipulation.",
        "    Peripheral Interfacing: GPIO, Timers, Interrupts, LCDs, Keypads.",
        "    State Machine Design: Implementing FSM for complex control logic.",
        "    Debugging: Using MPLAB X IDE and simulators."
    ])

    # Slide 4: Brief about Embedded Systems
    add_slide("Brief about Embedded Systems", [
        "Definition:",
        "    Computing system dedicated to a specific function within a larger system.",
        "Characteristics:",
        "    Specific Purpose: Designed for a singular task.",
        "    Resource Constrained: Limited RAM/Flash/Processing power.",
        "    Reliability: Expected to run continuously.",
        "    Real-Time: Strict timing deadlines (hard or soft real-time).",
        "Relevance:",
        "    Automotive, Consumer Electronics, Medical Devices, Industrial Automation."
    ])

    # Slide 5: Project Brief & Requirements
    add_slide("Project Brief & Requirements", [
        "Functional Requirements:",
        "    1. Power On: Welcome screen and initialization.",
        "    2. Menu System: Select between Micro, Grill, Convection, Start.",
        "    3. Input Handling: Matrix Keypad (0-9, *, #).",
        "    4. Display: Show mode, time, temperature on CLCD.",
        "    5. Load Control: Relay/Fan status (simulate heating).",
        "    6. Alerts: Buzzer notification on completion.",
        "Non-Functional Requirements:",
        "    Responsiveness, Accuracy (Timer), Usability."
    ])

    # Slide 6: System Design
    add_slide("System Design (Block Diagram)", [
        "Controller: PIC16F877A (20 MHz)",
        "Input: 4x4 Matrix Keypad (PORTD)",
        "    Rows and Columns scanning.",
        "Output: 16x4 CLCD (PORTD & PORTE)",
        "    Displays Menus, Timer, Temperature.",
        "Actuators:",
        "    Fan/Heater (RC2): Simulates Magnetron.",
        "    Buzzer (RC1): Audio feedback.",
        "Timer Module:",
        "    Timer2 used for 1-second interrupts."
    ])

    # Slide 7: Hardware & Simulation Environment
    add_slide("Hardware & Simulation Environment", [
        "Board/Simulator:",
        "    PICGenios Development Board / Proteus / PiximLab.",
        "    MCU: PIC16F877A (40-pin DIP).",
        "Key Hardware Specs:",
        "    GPIO: Digital I/O for Keypad and LCD.",
        "    Timers: Timer2 (8-bit) with Prescaler/Postscaler.",
        "    Interrupts: TMR2IE for non-blocking timing."
    ])

    # Slide 8: Development Process & Tools
    add_slide("Development Process & Tools", [
        "Software Tools:",
        "    MPLAB X IDE: Development Environment.",
        "    XC8 Compiler: C Compiler for PIC.",
        "    Simulators: For logic testing.",
        "Development Workflow:",
        "    1. Requirement Analysis.",
        "    2. Driver Development (LCD, Keypad).",
        "    3. Logic Implementation (FSM).",
        "    4. Integration.",
        "    5. Testing & Tuning (Timer calibration)."
    ])

    # Slide 9: Implementation Details - State Machine
    add_slide("Implementation Details - State Machine", [
        "Finite State Machine (FSM):",
        "    Manages system complexity and mode sets.",
        "States:",
        "    MAIN_MENU: Waiting for selection.",
        "    MICRO/GRILL/CONVECTION: Configuration.",
        "    TIME_DISPLAY: Active cooking state.",
        "    PAUSE/STOP: Control states.",
        "Logic:",
        "    Switch-case structure processing 'state' variable.",
        "    Transitions triggered by keypad input."
    ])

    # Slide 10: Implementation Details - Timer Logic
    add_slide("Implementation Details - Timer Logic", [
        "Challenge: keeping accurate time + responsive UI.",
        "Solution: Interrupt Service Routine (ISR).",
        "Mechanism:",
        "    Timer2 overflows every few ms.",
        "    ISR increments a counter.",
        "    At calibrated count, 'seconds' variable decrements.",
        "Calibration:",
        "    PR2=250, Prescaler=1:16.",
        "    Counter calibrated to 450 (empirically derived) for 1s accuracy."
    ])

    # Slide 11: Modular Project Structure
    add_slide("Modular Project Structure", [
        "main.c: Entry point, Super Loop, FSM.",
        "micro_oven.c: App specific screens/logic.",
        "isr.c: Interrupt handling (Timer).",
        "timers.c: Timer hardware config.",
        "clcd.c: LCD Driver.",
        "matrix_keypad.c: Keypad Driver."
    ])

    # Slide 12: Conclusion & Future Scope
    add_slide("Conclusion & Future Scope", [
        "Conclusion:",
        "    Successfully designed/simulated Microwave system.",
        "    Achieved accurate timing and responsive UI.",
        "    Demonstrated proficiency in Embedded C.",
        "Future Scope:",
        "    PWM Control for variable power.",
        "    EEPROM for saving custom programs.",
        "    Sensor integration (LM35) for temp control.",
        "    Hardware implementation on PCB."
    ])

    output_file = "Microwave_Project_Internship_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()
